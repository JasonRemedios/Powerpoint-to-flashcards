from pptx import Presentation
import pyperclip
import webbrowser

path = input('Enter path or name of powerpoint file: ')
try:
    prs = Presentation(path)

except:
    
    print('\nCould not load presentation. If you entered the name of the presentation, try entering the full path to the file.')
    path = input('Enter path or name of powerpoint file: ')
    prs = Presentation(path)
    
slides = prs.slides

cardList = []

for slide in slides:
    hasTitle = False
    hasBody = False
    #check for title and body and save index of each
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type == 1:
                hasTitle = True
                titleIndex = slide.shapes.index(shape)

            elif shape.placeholder_format.type == 2:
                hasBody = True
                bodyIndex = slide.shapes.index(shape)

    #format title and bullet points            
    if hasTitle == True and hasBody == True:
        slideStr = slide.shapes[titleIndex].text + '    '
        index = 1
        for i in slide.shapes[bodyIndex].text_frame.paragraphs:
            if i.level == 0:
                slideStr +=(' ' + str(index) + '.' + i.text)
                index += 1
            
            elif i.level == 1:
                slideStr += (' ' + '--' + i.text)

            elif i.level > 1:
                slideStr += (' ' + '**' + i.text)

        cardList.append(slideStr)

pyperclip.copy('\n'.join(cardList))
webbrowser.open('https://www.cram.com/flashcards/create')
